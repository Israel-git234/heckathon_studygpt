from flask import Flask, request, jsonify
from flask_cors import CORS
import os
import logging
from dotenv import load_dotenv
from config import Config
from services.course_builder import CourseBuilder
from functools import wraps
import traceback
from typing import Dict, Any, List
import time
from werkzeug.utils import secure_filename
from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

# Load environment variables
load_dotenv()

# Configure enhanced logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('app.log') if os.path.exists('.') else logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

app = Flask(__name__)

# Configure CORS with specific settings
CORS(app, 
     origins=[Config.FRONTEND_URL, "http://localhost:3000", "http://127.0.0.1:3000"],
     methods=['GET', 'POST', 'OPTIONS'],
     allow_headers=['Content-Type', 'Authorization']
)

# Configure Flask app
app.config['MAX_CONTENT_LENGTH'] = Config.MAX_CONTENT_LENGTH
app.config['JSON_SORT_KEYS'] = False

# Initialize course builder with error handling
try:
    course_builder = CourseBuilder()
    logger.info("Course builder initialized successfully")
except Exception as e:
    logger.error(f"Failed to initialize course builder: {e}")
    course_builder = None

# Validate configuration on startup
try:
    Config.validate_config()
    logger.info("Configuration validated successfully")
except ValueError as e:
    logger.error(f"Configuration error: {e}")
    logger.warning("Some features may not work without proper API keys")

# Error handler decorator
def handle_errors(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        try:
            return f(*args, **kwargs)
        except Exception as e:
            logger.error(f"Unhandled error in {f.__name__}: {e}")
            logger.error(traceback.format_exc())
            return jsonify({
                "error": "Internal server error",
                "message": "An unexpected error occurred. Please try again."
            }), 500
    return decorated_function

# Input validation helper
def validate_video_urls(data: Dict[str, Any]) -> List[str]:
    """Validate and extract video URLs from request data"""
    if not data:
        raise ValueError("No data provided")
    
    video_urls = data.get('video_urls', [])
    
    if not video_urls:
        raise ValueError("No video URLs provided")
    
    if not isinstance(video_urls, list):
        raise ValueError("video_urls must be a list")
    
    if len(video_urls) > 10:  # Reasonable limit
        raise ValueError("Too many videos. Maximum 10 videos allowed.")
    
    # Validate each URL
    validated_urls = []
    for url in video_urls:
        if not isinstance(url, str):
            raise ValueError("All video URLs must be strings")
        
        url = url.strip()
        if not url:
            continue
            
        if len(url) > 500:  # Reasonable URL length limit
            raise ValueError("URL too long")
            
        validated_urls.append(url)
    
    if not validated_urls:
        raise ValueError("No valid video URLs provided")
    
    return validated_urls

# Enhanced health check endpoint
@app.route('/api/health', methods=['GET'])
@handle_errors
def health_check():
    # Check service availability
    services_status = {
        "course_builder": course_builder is not None,
        "config_valid": True
    }
    
    try:
        Config.validate_config()
    except ValueError:
        services_status["config_valid"] = False
    
    status_code = 200 if all(services_status.values()) else 503
    
    return jsonify({
        "status": "healthy" if status_code == 200 else "degraded",
        "message": "StudyWeave AI Backend Running",
        "services": services_status,
        "version": "1.0.0"
    }), status_code

# Main course generation endpoint
@app.route('/api/generate-course', methods=['POST'])
@handle_errors
def generate_course():
    """
    Main endpoint for generating structured courses from YouTube videos
    Processes video URLs and returns structured course data with concepts and quizzes
    """
    # Check if course builder is available
    if not course_builder:
        return jsonify({
            "error": "Service unavailable",
            "message": "Course generation service is not available"
        }), 503
    
    # Get and validate request data
    data = request.get_json(force=True)
    if not data:
        return jsonify({
            "error": "Invalid request",
            "message": "No JSON data provided"
        }), 400
    
    try:
        video_urls = validate_video_urls(data)
    except ValueError as e:
        return jsonify({
            "error": "Invalid input",
            "message": str(e)
        }), 400
    
    logger.info(f"Generating course from {len(video_urls)} video URLs")
    
    # Build the course using our services
    course_data = course_builder.build_course_from_videos(video_urls)
    
    if isinstance(course_data, dict) and "error" in course_data:
        return jsonify(course_data), 400
    
    # Add metadata
    course_data["generated_at"] = int(time.time())
    course_data["api_version"] = "1.0.0"
    
    logger.info(f"Successfully generated course with {course_data.get('total_concepts', 0)} concepts")
    return jsonify(course_data)

# Preview endpoint for quick video info
@app.route('/api/preview-videos', methods=['POST'])
@handle_errors
def preview_videos():
    """
    Quick preview of video metadata without AI processing
    Useful for validating URLs before full course generation
    """
    # Check if course builder is available
    if not course_builder:
        return jsonify({
            "error": "Service unavailable",
            "message": "Video preview service is not available"
        }), 503
    
    # Get and validate request data
    data = request.get_json(force=True)
    if not data:
        return jsonify({
            "error": "Invalid request",
            "message": "No JSON data provided"
        }), 400
    
    try:
        video_urls = validate_video_urls(data)
    except ValueError as e:
        return jsonify({
            "error": "Invalid input",
            "message": str(e)
        }), 400
    
    logger.info(f"Previewing {len(video_urls)} video URLs")
    
    # Get just video info without AI processing
    preview_data = course_builder.get_video_info_only(video_urls)
    
    # Add metadata
    preview_data["previewed_at"] = int(time.time())
    preview_data["api_version"] = "1.0.0"
    
    return jsonify(preview_data)

# Ask-AI endpoint
@app.route('/api/ask-question', methods=['POST'])
@handle_errors
def ask_question():
    if not course_builder:
        return jsonify({"error": "Service unavailable", "message": "AI service not available"}), 503

    data = request.get_json(force=True)
    if not data or not isinstance(data.get('question'), str):
        return jsonify({"error": "Invalid input", "message": "'question' is required"}), 400

    video = data.get('video') or {}
    concept = data.get('concept') or {}

    answer = course_builder.ai_service.answer_question(
        question=data['question'],
        video_data=video if isinstance(video, dict) else None,
        concept=concept if isinstance(concept, dict) else None,
    )

    return jsonify({
        "answer": answer,
        "answered_at": int(time.time()),
        "api_version": "1.0.0"
    })

# Helpers to extract text from uploads
def extract_text_from_pdf(stream: BytesIO) -> str:
    try:
        reader = PdfReader(stream)
        return "\n".join(page.extract_text() or '' for page in reader.pages)
    except Exception:
        return ''

def extract_text_from_docx(stream: BytesIO) -> str:
    try:
        doc = DocxDocument(stream)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ''

def extract_text_from_pptx(stream: BytesIO) -> str:
    try:
        prs = Presentation(stream)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ''

ALLOWED_EXT = {'.pdf', '.docx', '.pptx'}

@app.route('/api/summarize-upload', methods=['POST'])
@handle_errors
def summarize_upload():
    """Accept a file (PDF/DOCX/PPTX) or raw transcript text and return study notes + suggested videos."""
    text = None
    if 'file' in request.files:
        f = request.files['file']
        filename = secure_filename(f.filename or '')
        ext = ('.' + filename.rsplit('.', 1)[-1].lower()) if '.' in filename else ''
        data = BytesIO(f.read())
        if ext == '.pdf':
            text = extract_text_from_pdf(data)
        elif ext == '.docx':
            text = extract_text_from_docx(data)
        elif ext == '.pptx':
            text = extract_text_from_pptx(data)
        else:
            return jsonify({"error": "Invalid file type", "message": "Use PDF, PPTX, or DOCX"}), 400
    else:
        payload = request.get_json(silent=True) or {}
        text = payload.get('transcript')

    if not text or len(text.strip()) < 30:
        return jsonify({"error": "Invalid input", "message": "Provide a valid file or transcript text"}), 400

    # Ask AI to produce structured study notes
    ai = course_builder.ai_service
    prompt = (
        "You are a helpful educator. Given lecture text/transcript, produce concise study notes with: "
        "summary (4-6 sentences), 5 key bullet points, 3 terminology definitions, and 2 practice questions.\n\n"
        f"TEXT:\n{text[:12000]}"
    )
    try:
        resp = ai.model.generate_content(prompt, generation_config=ai.genai.types.GenerationConfig if hasattr(ai, 'genai') else None)
        notes = resp.text.strip() if hasattr(resp, 'text') else ''
    except Exception:
        notes = ""

    # fallback minimal notes
    if not notes:
        notes = "Summary: Lecture overview unavailable. Key points could not be fully extracted."

    # Suggest videos by topic using first 10 words as query
    topic = " ".join(text.strip().split()[:10])
    yt = course_builder.youtube_service
    # We will return placeholder suggestions from our processed titles if API quota limits; frontend can still preview
    suggestions = []
    try:
        # Use YouTube Data API search
        search = yt.youtube.search().list(part='snippet', q=topic, type='video', maxResults=5)
        res = search.execute()
        for item in res.get('items', []):
            vid = item['id']['videoId']
            snippet = item['snippet']
            suggestions.append({
                'id': vid,
                'title': snippet.get('title',''),
                'channel': snippet.get('channelTitle',''),
                'thumbnail': snippet.get('thumbnails',{}).get('high',{}).get('url',''),
                'url': f'https://www.youtube.com/watch?v={vid}'
            })
    except Exception:
        suggestions = []

    return jsonify({
        'notes': notes,
        'recommended_videos': suggestions,
        'api_version': '1.0.0',
        'generated_at': int(time.time())
    })

# Global error handlers
@app.errorhandler(404)
def not_found(error):
    return jsonify({
        "error": "Not found",
        "message": "The requested endpoint does not exist"
    }), 404

@app.errorhandler(405)
def method_not_allowed(error):
    return jsonify({
        "error": "Method not allowed",
        "message": "The requested method is not allowed for this endpoint"
    }), 405

@app.errorhandler(413)
def request_entity_too_large(error):
    return jsonify({
        "error": "Request too large",
        "message": "The request payload is too large"
    }), 413

if __name__ == '__main__':
    import time
    app.run(debug=Config.DEBUG, port=5000, host='127.0.0.1')
